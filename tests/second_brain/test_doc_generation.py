"""Tests for second_brain.doc_generation -- the minimal markdown-driven
docx/pptx pipeline. Covers the parser and both renderers structurally
(real headings/bullets/tables/slides come out, not just "a file exists");
does not cover to_pdf() (requires the live soffice binary, environment-
dependent, not something to gate CI-less local test runs on)."""
from docx import Document
from pptx import Presentation

from second_brain.doc_generation import (
    _parse_blocks,
    _parse_inline_runs,
    markdown_to_docx,
    markdown_to_pptx,
)

SAMPLE_MD = """\
# Report Title

## Section One

Some intro text with **bold** and *italic* words.

- first bullet
- second bullet

## Section Two

| A | B |
|---|---|
| 1 | 2 |
| 3 | 4 |

A closing paragraph.
"""


def test_parse_blocks_shape():
    blocks = _parse_blocks(SAMPLE_MD)
    types = [b["type"] for b in blocks]
    assert types == [
        "heading", "heading", "paragraph", "bullet", "bullet",
        "heading", "table", "paragraph",
    ]
    assert blocks[0] == {"type": "heading", "level": 1, "text": "Report Title"}
    assert blocks[5]["text"] == "Section Two"
    assert blocks[6]["rows"] == [["A", "B"], ["1", "2"], ["3", "4"]]


def test_parse_inline_runs_bold_and_italic():
    runs = _parse_inline_runs("plain **bold** mid *italic* end")
    assert ("plain ", False, False) in runs
    assert ("bold", True, False) in runs
    assert (" mid ", False, False) in runs
    assert ("italic", False, True) in runs


def test_markdown_to_docx_structure(tmp_path):
    out = markdown_to_docx(SAMPLE_MD, str(tmp_path / "out.docx"), title="Test Doc")
    doc = Document(out)
    heading_styles = [p.style.name for p in doc.paragraphs if p.style.name.startswith("Heading") or p.style.name == "Title"]
    assert "Title" in heading_styles
    assert "Heading 1" in heading_styles
    assert "Heading 2" in heading_styles
    assert len(doc.tables) == 1
    table = doc.tables[0]
    assert [c.text for c in table.rows[0].cells] == ["A", "B"]
    assert [c.text for c in table.rows[1].cells] == ["1", "2"]
    bullet_paras = [p for p in doc.paragraphs if p.style.name == "List Bullet"]
    assert len(bullet_paras) == 2
    assert bullet_paras[0].text == "first bullet"


def test_markdown_to_pptx_outline_slides(tmp_path):
    out = markdown_to_pptx(SAMPLE_MD, str(tmp_path / "out.pptx"), title="Test Deck")
    prs = Presentation(out)
    # Title slide + one slide per H1/H2 (Report Title, Section One, Section Two) = 4
    assert len(prs.slides) == 4
    titles = [s.shapes.title.text if s.shapes.title else "" for s in prs.slides]
    assert titles[0] == "Test Deck"
    assert "Section One" in titles
    assert "Section Two" in titles


def test_markdown_to_docx_empty_input_does_not_crash(tmp_path):
    out = markdown_to_docx("", str(tmp_path / "empty.docx"), title="Empty")
    doc = Document(out)
    assert doc.paragraphs[0].text == "Empty"
