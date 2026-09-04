"""
second_brain.doc_generation -- turn vault content (markdown) into real
documents: .docx, .pptx, and PDF (via LibreOffice headless conversion).

General-purpose second-brain tooling, not investor-doc-specific: any vault
note or synthesis (a 04-Syntheses/ digest, a research write-up, a weekly
compile) can go through this to become a real handoff-able document instead
of staying markdown-only. The investor-materials use case (2026-08-24) is
the first caller, not the reason this exists.

Deliberately minimal for a first pass -- a real markdown parser (headings,
bullet lists, simple tables, bold/italic inline runs) mapped onto
python-docx/python-pptx primitives, not a full CommonMark implementation.
Nested lists, code blocks, images, and footnotes are NOT handled; call
render_docx()/render_pptx() directly with a structured spec instead of
going through the markdown parser if you need something the parser can't
express. A more polished version (real slide templates/branding, chart
embedding) is planned as a follow-up once this minimal version has proven
out the content-mapping approach -- see the operator conversation this
was built from, 2026-08-24.

Three entry points:
  markdown_to_docx(markdown_text, output_path, title=None) -> path
  markdown_to_pptx(markdown_text, output_path, title=None) -> path
      (H1/H2 headings start a new slide; bullet lines under a heading
      become that slide's bullets -- the common "outline as deck"
      convention, not a full markdown-to-slides heuristic.)
  to_pdf(path) -> pdf_path
      (LibreOffice headless conversion -- works on the .docx/.pptx this
      module produces, or any other docx/pptx/odt/etc. LibreOffice opens.)

CLI:
  python3 -m second_brain.doc_generation docx input.md output.docx
  python3 -m second_brain.doc_generation pptx input.md output.pptx
  python3 -m second_brain.doc_generation pdf output.docx
"""
from __future__ import annotations

import re
import subprocess
import sys
from pathlib import Path

from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.ns import qn
from docx.shared import Pt, RGBColor, Inches
from pptx import Presentation
from pptx.dml.color import RGBColor as PptxRGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt as PptxPt, Inches as PptxInches, Emu

# ---------------------------------------------------------------------------
# CTDI brand theme -- 2026-09-04, extracted directly (not eyeballed) from
# the v1.1 investor-materials set (docs/investor-materials/v1.1/), which
# the operator set as the design benchmark after v1.2-draft/v1.5's output
# (the plain default-python-docx/pptx look below) was rejected as "generic
# nonsense". Colors pulled from the actual .docx/.pptx XML (run.font.color,
# table cell shading, shape fill/line) via python-docx/python-pptx
# introspection, not estimated from a rendered screenshot -- see the
# session's own extraction commands if this ever needs re-deriving.
# ---------------------------------------------------------------------------
CTDI_THEME = {
    # pptx dark-navy brand system
    "navy_bg": "0A1628",
    "card_fill": "111F33",
    "badge_fill": "16283F",
    "gold": "E8A124",
    "white": "FFFFFF",
    "muted": "8FA3BF",
    "green": "6FBF8F",
    # docx light-paper brand system (same brand, print-contrast variant)
    "gold_docx": "B57A0E",
    "title_docx": "0A1628",
    "subtitle_docx": "444E5A",
    "byline_docx": "5A626B",
    "body_docx": "1A222E",
    "table_header_docx": "0A1628",
    "table_zebra_docx": "F2F5F9",
    "green_docx": "1E6E3C",
    "blue_docx": "145A96",
    "gray_docx": "5A626B",
    # fonts -- Montserrat may not be installed wherever a file is later
    # opened/converted; declaring it is still correct (most investor-
    # facing readers open in real Office/Google Docs, which do have it),
    # and every renderer here falls back gracefully if it's absent.
    "heading_font": "Montserrat",
    "body_font": "Calibri",
}

# Tier-status strings this repo's investor docs use verbatim (see
# docs/investor-materials/*/src/executive-summary.md's "Evidence
# discipline" convention) -- auto-colored wherever they appear as a bold
# run, in either docx or pptx, so a table/bullet doesn't need markup
# beyond the **bold** the source markdown already uses.
_TIER_COLORS_DOCX = {
    "LIVE & VERIFIED": "1E6E3C",
    "CODE-COMPLETE, NOT LIVE-VERIFIED": "145A96",
    "CODE-COMPLETE": "145A96",
    "ROADMAP": "5A626B",
}
_TIER_COLORS_PPTX = {
    "LIVE & VERIFIED": "6FBF8F",
    "CODE-COMPLETE, NOT LIVE-VERIFIED": "6FA8DC",
    "CODE-COMPLETE": "6FA8DC",
    "ROADMAP": "8FA3BF",
}

# 2026-09-04: the actual v1.5 source markdown (docs/investor-materials/v1.5/
# */src/*.md) writes these tier labels as plain leading text in a table
# cell -- "LIVE and VERIFIED", not "**LIVE & VERIFIED**" -- so detecting
# them only on bold runs (which is how v1.1's own hand-built docx/pptx
# apparently did it) silently colors nothing. Matched at the START of a
# cell/run's text, case-insensitive, "and"/"&" both accepted; only the
# matched tier phrase itself is colored, trailing explanatory text after
# it (", with the availability caveat below") stays normal body color.
_TIER_PATTERN = re.compile(
    r"^(LIVE\s+(?:and|&)\s+VERIFIED|CODE-COMPLETE|ROADMAP)\b", re.IGNORECASE
)


def _split_tier_prefix(text: str) -> tuple[str, str] | None:
    """Returns (tier_phrase, rest_of_text) if text starts with a known
    tier label, else None."""
    m = _TIER_PATTERN.match(text.strip())
    if not m:
        return None
    tier_phrase = m.group(1)
    rest = text.strip()[m.end():]
    return tier_phrase, rest


# ---------------------------------------------------------------------------
# Minimal markdown parsing -- headings, bullets, tables, plain paragraphs.
# Inline bold (**x**) and italic (*x*) are handled per-run within a
# paragraph; nothing else inline (links render as literal text, which is
# an acceptable minimal-pass tradeoff -- most vault content doesn't lean
# on markdown links for meaning).
# ---------------------------------------------------------------------------

_HEADING_RE = re.compile(r"^(#{1,6})\s+(.*)$")
_BULLET_RE = re.compile(r"^\s*[-*]\s+(.*)$")
_TABLE_ROW_RE = re.compile(r"^\s*\|(.+)\|\s*$")
_TABLE_SEP_RE = re.compile(r"^\s*\|[\s:|-]+\|\s*$")


def _parse_inline_runs(text: str) -> list[tuple[str, bool, bool]]:
    """Split a line into (text, bold, italic) runs on **bold** / *italic*
    markers. Not a real tokenizer -- doesn't handle nested or overlapping
    emphasis, which is fine for the vault content this targets."""
    runs: list[tuple[str, bool, bool]] = []
    pos = 0
    pattern = re.compile(r"\*\*(.+?)\*\*|\*(.+?)\*")
    for m in pattern.finditer(text):
        if m.start() > pos:
            runs.append((text[pos:m.start()], False, False))
        if m.group(1) is not None:
            runs.append((m.group(1), True, False))
        else:
            runs.append((m.group(2), False, True))
        pos = m.end()
    if pos < len(text):
        runs.append((text[pos:], False, False))
    return runs or [(text, False, False)]


def _parse_blocks(markdown_text: str) -> list[dict]:
    """Parse markdown into a flat list of block dicts:
    {"type": "heading", "level": int, "text": str}
    {"type": "bullet", "text": str}
    {"type": "table", "rows": list[list[str]]}
    {"type": "paragraph", "text": str}
    Blank lines are separators only, not emitted as blocks."""
    blocks: list[dict] = []
    lines = markdown_text.splitlines()
    i = 0
    while i < len(lines):
        line = lines[i]
        if not line.strip():
            i += 1
            continue
        m = _HEADING_RE.match(line)
        if m:
            blocks.append({"type": "heading", "level": len(m.group(1)), "text": m.group(2).strip()})
            i += 1
            continue
        m = _BULLET_RE.match(line)
        if m:
            blocks.append({"type": "bullet", "text": m.group(1).strip()})
            i += 1
            continue
        if _TABLE_ROW_RE.match(line):
            rows = []
            while i < len(lines) and _TABLE_ROW_RE.match(lines[i]):
                if not _TABLE_SEP_RE.match(lines[i]):
                    cells = [c.strip() for c in lines[i].strip().strip("|").split("|")]
                    rows.append(cells)
                i += 1
            if rows:
                blocks.append({"type": "table", "rows": rows})
            continue
        # Plain paragraph -- collect contiguous non-blank, non-special lines.
        para_lines = [line.strip()]
        i += 1
        while i < len(lines) and lines[i].strip() and not _HEADING_RE.match(lines[i]) \
                and not _BULLET_RE.match(lines[i]) and not _TABLE_ROW_RE.match(lines[i]):
            para_lines.append(lines[i].strip())
            i += 1
        blocks.append({"type": "paragraph", "text": " ".join(para_lines)})
    return blocks


# ---------------------------------------------------------------------------
# docx rendering
# ---------------------------------------------------------------------------

def render_docx(blocks: list[dict], output_path: str, title: str | None = None) -> str:
    """Render parsed blocks (see _parse_blocks) into a .docx file."""
    doc = Document()
    if title:
        doc.add_heading(title, level=0)

    for block in blocks:
        if block["type"] == "heading":
            doc.add_heading(block["text"], level=min(block["level"], 4))
        elif block["type"] == "bullet":
            p = doc.add_paragraph(style="List Bullet")
            for text, bold, italic in _parse_inline_runs(block["text"]):
                run = p.add_run(text)
                run.bold = bold
                run.italic = italic
        elif block["type"] == "table":
            rows = block["rows"]
            if not rows:
                continue
            table = doc.add_table(rows=len(rows), cols=len(rows[0]))
            table.style = "Light Grid Accent 1"
            for r, row in enumerate(rows):
                for c, cell_text in enumerate(row):
                    if c < len(table.rows[r].cells):
                        table.rows[r].cells[c].text = cell_text
        elif block["type"] == "paragraph":
            p = doc.add_paragraph()
            for text, bold, italic in _parse_inline_runs(block["text"]):
                run = p.add_run(text)
                run.bold = bold
                run.italic = italic
                run.font.size = Pt(11)

    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    doc.save(output_path)
    return output_path


def markdown_to_docx(markdown_text: str, output_path: str, title: str | None = None) -> str:
    return render_docx(_parse_blocks(markdown_text), output_path, title=title)


def _rgb(hex_str: str) -> RGBColor:
    return RGBColor.from_string(hex_str)


def _shade_cell(cell, hex_fill: str) -> None:
    """Set a table cell's background shading (python-docx has no public
    API for this -- direct w:shd XML element, same technique Word itself
    uses)."""
    tcPr = cell._tc.get_or_add_tcPr()
    shd = tcPr.find(qn("w:shd"))
    if shd is None:
        shd = tcPr.makeelement(qn("w:shd"), {})
        tcPr.append(shd)
    shd.set(qn("w:val"), "clear")
    shd.set(qn("w:color"), "auto")
    shd.set(qn("w:fill"), hex_fill)


def _tier_docx_color(text: str) -> str | None:
    key = re.sub(r"\bAND\b", "&", text.strip().upper())
    for tier, color in _TIER_COLORS_DOCX.items():
        if key.startswith(tier):
            return color
    return None


def _add_inline_runs_docx(paragraph, text: str, theme: dict, base_color: str | None = None,
                           size: int | None = None) -> None:
    for run_text, bold, italic in _parse_inline_runs(text):
        run = paragraph.add_run(run_text)
        run.bold = bold
        run.italic = italic
        if size:
            run.font.size = Pt(size)
        run.font.name = theme["body_font"]
        tier_color = _tier_docx_color(run_text) if bold else None
        if tier_color:
            run.font.color.rgb = _rgb(tier_color)
        elif base_color:
            run.font.color.rgb = _rgb(base_color)


def render_docx_branded(blocks: list[dict], output_path: str, theme: dict = CTDI_THEME) -> str:
    """Branded docx render matching the v1.1 design benchmark (see
    CTDI_THEME's own comment for where the palette came from): a styled
    title block (eyebrow-less on docx -- v1.1's docx eyebrow line IS the
    first block, handled below), navy section headings, and tables with
    a navy header row, white bold header text, and light zebra striping,
    with LIVE & VERIFIED / CODE-COMPLETE / ROADMAP tier strings
    auto-colored wherever they appear as a bold run.

    Expects the source markdown's own established convention (see any
    docs/investor-materials/*/src/*.md): the first 1-3 non-heading blocks
    before the first heading are, in order, an eyebrow/title line, an
    italic subtitle, and a byline -- exactly what
    docs/investor-materials/v1.5/*/src/{executive-summary,due-diligence-faq}.md
    already write. Falls back gracefully (plain paragraphs) if fewer than
    3 leading blocks are present."""
    doc = Document()
    section = doc.sections[0]
    section.left_margin = section.right_margin = Inches(1)

    # ---- Title block (first up-to-3 leading paragraph blocks) ----
    lead = []
    body_start = 0
    for b in blocks:
        if b["type"] != "paragraph" or len(lead) >= 3:
            break
        lead.append(b["text"])
        body_start += 1

    if lead:
        eyebrow_or_title = lead[0]
        p = doc.add_paragraph()
        run = p.add_run(eyebrow_or_title)
        run.font.name = theme["heading_font"]
        run.font.size = Pt(24)
        run.bold = True
        run.font.color.rgb = _rgb(theme["title_docx"])

    if len(lead) > 1:
        p = doc.add_paragraph()
        run = p.add_run(lead[1])
        run.font.name = theme["body_font"]
        run.font.size = Pt(11)
        run.italic = True
        run.font.color.rgb = _rgb(theme["subtitle_docx"])

    if len(lead) > 2:
        p = doc.add_paragraph()
        run = p.add_run(lead[2])
        run.font.name = theme["body_font"]
        run.font.size = Pt(9)
        run.font.color.rgb = _rgb(theme["byline_docx"])

    if lead:
        # Gold horizontal rule -- bottom border on an empty paragraph,
        # same technique Word's own "Horizontal Line" feature uses.
        p = doc.add_paragraph()
        pPr = p._p.get_or_add_pPr()
        pBdr = pPr.makeelement(qn("w:pBdr"), {})
        bottom = pPr.makeelement(qn("w:bottom"), {
            qn("w:val"): "single", qn("w:sz"): "18", qn("w:space"): "1",
            qn("w:color"): theme["gold_docx"],
        })
        pBdr.append(bottom)
        pPr.append(pBdr)

    # ---- Body ----
    for block in blocks[body_start:]:
        if block["type"] == "heading":
            p = doc.add_paragraph()
            run = p.add_run(block["text"])
            run.font.name = theme["heading_font"]
            run.bold = True
            run.font.size = Pt(16 if block["level"] <= 1 else 13)
            run.font.color.rgb = _rgb(theme["title_docx"])
            p.paragraph_format.space_before = Pt(14)
            p.paragraph_format.space_after = Pt(6)
        elif block["type"] == "bullet":
            p = doc.add_paragraph(style="List Bullet")
            _add_inline_runs_docx(p, block["text"], theme, base_color=theme["body_docx"], size=11)
        elif block["type"] == "table":
            rows = block["rows"]
            if not rows:
                continue
            table = doc.add_table(rows=len(rows), cols=len(rows[0]))
            table.style = "Table Grid"
            for r, row in enumerate(rows):
                for c, cell_text in enumerate(row):
                    if c >= len(table.rows[r].cells):
                        continue
                    cell = table.rows[r].cells[c]
                    cell.text = ""
                    p = cell.paragraphs[0]
                    if r == 0:
                        _shade_cell(cell, theme["table_header_docx"])
                        run = p.add_run(cell_text)
                        run.bold = True
                        run.font.name = theme["body_font"]
                        run.font.size = Pt(10)
                        run.font.color.rgb = _rgb(theme["white"])
                    else:
                        if r % 2 == 0:
                            _shade_cell(cell, theme["table_zebra_docx"])
                        tier_split = _split_tier_prefix(cell_text)
                        if tier_split:
                            tier_phrase, rest = tier_split
                            tier_run = p.add_run(tier_phrase)
                            tier_run.bold = True
                            tier_run.font.name = theme["body_font"]
                            tier_run.font.size = Pt(10)
                            tier_run.font.color.rgb = _rgb(_tier_docx_color(tier_phrase))
                            if rest:
                                _add_inline_runs_docx(p, rest, theme, base_color=theme["body_docx"], size=10)
                        else:
                            _add_inline_runs_docx(p, cell_text, theme, base_color=theme["body_docx"], size=10)
        elif block["type"] == "paragraph":
            p = doc.add_paragraph()
            p.paragraph_format.space_after = Pt(8)
            _add_inline_runs_docx(p, block["text"], theme, base_color=theme["body_docx"], size=11)

    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    doc.save(output_path)
    return output_path


def markdown_to_docx_branded(markdown_text: str, output_path: str, theme: dict = CTDI_THEME) -> str:
    return render_docx_branded(_parse_blocks(markdown_text), output_path, theme=theme)


# ---------------------------------------------------------------------------
# pptx rendering -- outline convention: each H1/H2 starts a new slide;
# bullets and paragraphs under it become that slide's body content.
# ---------------------------------------------------------------------------

def render_pptx(slides: list[dict], output_path: str, title: str | None = None) -> str:
    """slides: list of {"title": str, "bullets": list[str]}.
    Uses the default template's Title Slide + Title-and-Content layouts --
    no custom branding/theme in this minimal pass (see module docstring)."""
    prs = Presentation()

    if title:
        title_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_layout)
        slide.shapes.title.text = title

    content_layout = prs.slide_layouts[1]
    for spec in slides:
        slide = prs.slides.add_slide(content_layout)
        slide.shapes.title.text = spec.get("title", "")
        body = slide.placeholders[1].text_frame
        bullets = spec.get("bullets", [])
        if not bullets:
            continue
        body.text = bullets[0]
        body.paragraphs[0].font.size = PptxPt(18)
        for b in bullets[1:]:
            p = body.add_paragraph()
            p.text = b
            p.font.size = PptxPt(18)

    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    return output_path


def _blocks_to_slides(blocks: list[dict]) -> list[dict]:
    """Outline convention: an H1/H2 starts a new slide; bullets and
    paragraphs accumulate as that slide's bullet points until the next
    H1/H2. H3+ headings are demoted to a bold-leading bullet rather than
    a new slide (a deck this shallow shouldn't nest that deep)."""
    slides: list[dict] = []
    current: dict | None = None
    for block in blocks:
        if block["type"] == "heading" and block["level"] <= 2:
            current = {"title": block["text"], "bullets": []}
            slides.append(current)
            continue
        if current is None:
            current = {"title": "", "bullets": []}
            slides.append(current)
        if block["type"] == "bullet":
            current["bullets"].append(block["text"])
        elif block["type"] == "heading":
            current["bullets"].append(block["text"])
        elif block["type"] == "paragraph":
            current["bullets"].append(block["text"])
        elif block["type"] == "table":
            for row in block["rows"]:
                current["bullets"].append(" | ".join(row))
    return slides


def markdown_to_pptx(markdown_text: str, output_path: str, title: str | None = None) -> str:
    blocks = _parse_blocks(markdown_text)
    slides = _blocks_to_slides(blocks)
    return render_pptx(slides, output_path, title=title)


# ---------------------------------------------------------------------------
# Branded pptx rendering -- CTDI_THEME's dark-navy/gold system, matching
# the v1.1 design benchmark. Card-per-bullet grid (not the exact title+body
# 2-part cards v1.1 hand-built for its 4-item "problem" slide -- this
# repo's source markdown writes single flowing bullet sentences, not a
# title/body pair per bullet, so inventing a card title by truncating text
# would be a fabrication, not a faithful re-derivation; a card per full
# bullet keeps every word the source markdown actually wrote) for slides
# with a small, grid-friendly bullet count; a clean bulleted list
# (gold dot markers, on the same navy background) otherwise -- long lists
# and table-derived rows don't read well as cards.
# ---------------------------------------------------------------------------

_CARD_GRID_MAX_BULLETS = 6


def _pptx_rgb(hex_str: str) -> PptxRGBColor:
    return PptxRGBColor.from_string(hex_str)


def _tier_pptx_color(text: str) -> str | None:
    key = re.sub(r"\bAND\b", "&", text.strip().upper())
    for tier, color in _TIER_COLORS_PPTX.items():
        if key.startswith(tier):
            return color
    return None


def _set_slide_background(slide, hex_color: str) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _pptx_rgb(hex_color)


def _add_textbox(slide, left, top, width, height, text: str, *, size: int, bold: bool = False,
                  italic: bool = False, color: str, font: str, align=PP_ALIGN.LEFT,
                  tier_aware: bool = False):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    if tier_aware:
        for run_text, rbold, ritalic in _parse_inline_runs(text):
            run = p.add_run()
            run.text = run_text
            run.font.size = PptxPt(size)
            run.font.bold = bold or rbold
            run.font.italic = italic or ritalic
            run.font.name = font
            tc = _tier_pptx_color(run_text) if rbold else None
            run.font.color.rgb = _pptx_rgb(tc or color)
    else:
        run = p.add_run()
        run.text = text
        run.font.size = PptxPt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.name = font
        run.font.color.rgb = _pptx_rgb(color)
    return box


def _add_footer(slide, theme: dict, page_num: int | None, footer_text: str) -> None:
    _add_textbox(slide, PptxInches(0.6), PptxInches(7.08), PptxInches(9.5), PptxInches(0.3),
                 footer_text, size=9, color=theme["muted"], font=theme["body_font"])
    if page_num is not None:
        _add_textbox(slide, PptxInches(12.45), PptxInches(7.08), PptxInches(0.55), PptxInches(0.3),
                     str(page_num), size=9, color=theme["muted"], font=theme["body_font"],
                     align=PP_ALIGN.RIGHT)


def _add_bullet_cards(slide, bullets: list[str], theme: dict, top_emu: int) -> None:
    """2-column card grid, one full bullet per card -- see this section's
    module comment for why it's one bullet per card, not a title+body
    split the source markdown doesn't provide."""
    n = len(bullets)
    cols = 2 if n > 1 else 1
    rows = (n + cols - 1) // cols
    margin = PptxInches(0.6)
    gap = PptxInches(0.28)
    usable_w = PptxInches(13.333) - 2 * margin
    card_w = int((usable_w - gap * (cols - 1)) / cols)
    card_h = PptxInches(1.15) if rows <= 2 else PptxInches(0.95)
    for idx, bullet in enumerate(bullets):
        col, row = idx % cols, idx // cols
        left = margin + col * (card_w + gap)
        top = top_emu + row * (card_h + gap)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, card_w, card_h)
        card.adjustments[0] = 0.06
        card.fill.solid()
        card.fill.fore_color.rgb = _pptx_rgb(theme["card_fill"])
        card.line.fill.background()
        card.shadow.inherit = False
        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = PptxInches(0.22)
        tf.margin_top = tf.margin_bottom = PptxInches(0.14)
        tf.vertical_anchor = 1  # MSO_ANCHOR.MIDDLE -- keeps short/long cards visually centered
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        for run_text, rbold, ritalic in _parse_inline_runs(bullet):
            run = p.add_run()
            run.text = run_text
            run.font.size = PptxPt(13)
            run.font.bold = rbold
            run.font.italic = ritalic
            run.font.name = theme["body_font"]
            tc = _tier_pptx_color(run_text) if rbold else None
            run.font.color.rgb = _pptx_rgb(tc or theme["white"])


def _add_bullet_list(slide, bullets: list[str], theme: dict, top_emu: int) -> None:
    box = slide.shapes.add_textbox(PptxInches(0.6), Emu(top_emu), PptxInches(12.1), PptxInches(5.2))
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for bullet in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = PptxPt(10)
        marker = p.add_run()
        marker.text = "●  "
        marker.font.size = PptxPt(12)
        marker.font.color.rgb = _pptx_rgb(theme["gold"])
        marker.font.name = theme["body_font"]
        for run_text, rbold, ritalic in _parse_inline_runs(bullet):
            run = p.add_run()
            run.text = run_text
            run.font.size = PptxPt(14)
            run.font.bold = rbold
            run.font.italic = ritalic
            run.font.name = theme["body_font"]
            tc = _tier_pptx_color(run_text) if rbold else None
            run.font.color.rgb = _pptx_rgb(tc or theme["white"])


def render_pptx_branded(slides: list[dict], output_path: str, footer_text: str = "",
                         theme: dict = CTDI_THEME) -> str:
    """Branded pptx render matching the v1.1 design benchmark: solid navy
    background on every slide, gold-accented white-bold titles, a
    per-bullet card grid for short bullet lists (see _add_bullet_cards),
    a plain gold-marker bulleted list for longer/table-derived content,
    and a muted footer + page number on every content slide. First entry
    in `slides` is rendered as the title slide (no footer/page number,
    larger centered-ish title block) -- matches _blocks_to_slides()'s
    convention where the deck's first H2 IS the title block."""
    prs = Presentation()
    prs.slide_width = PptxInches(13.333)
    prs.slide_height = PptxInches(7.5)
    blank_layout = prs.slide_layouts[6]  # fully blank layout, no placeholders to fight

    for idx, spec in enumerate(slides):
        slide = prs.slides.add_slide(blank_layout)
        _set_slide_background(slide, theme["navy_bg"])
        title_text = spec.get("title", "")
        bullets = spec.get("bullets", [])

        if idx == 0:
            # Title slide. Long titles (e.g. "CTDI -- Corporate Travel
            # Dispatch Intelligence") wrap to 2 lines at 44pt and need a
            # taller box + a lower subtitle start, or the two blocks
            # visually collide -- shrink + drop for anything past ~28
            # chars rather than clipping/overlapping.
            title_size = 44 if len(title_text) <= 28 else 36
            title_box_h = PptxInches(1.3) if len(title_text) <= 28 else PptxInches(2.0)
            subtitle_top = PptxInches(3.6) if len(title_text) <= 28 else PptxInches(4.3)
            _add_textbox(slide, PptxInches(0.8), PptxInches(2.2), PptxInches(11.5), title_box_h,
                         title_text, size=title_size, bold=True, color=theme["white"], font=theme["heading_font"])
            if bullets:
                _add_textbox(slide, PptxInches(0.8), subtitle_top, PptxInches(10.8), PptxInches(2.6),
                             "\n".join(bullets), size=15, color=theme["muted"], font=theme["body_font"])
            if footer_text:
                _add_textbox(slide, PptxInches(0.8), PptxInches(7.05), PptxInches(10), PptxInches(0.3),
                             footer_text, size=10, color=theme["muted"], font=theme["body_font"])
            continue

        _add_textbox(slide, PptxInches(0.6), PptxInches(0.35), PptxInches(12.1), PptxInches(0.85),
                     title_text, size=30, bold=True, color=theme["white"], font=theme["heading_font"])
        content_top = int(PptxInches(1.5))
        if not bullets:
            pass
        elif len(bullets) <= _CARD_GRID_MAX_BULLETS and all(len(b) < 220 for b in bullets):
            _add_bullet_cards(slide, bullets, theme, content_top)
        else:
            _add_bullet_list(slide, bullets, theme, content_top)
        _add_footer(slide, theme, idx + 1, footer_text)

    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    return output_path


def markdown_to_pptx_branded(markdown_text: str, output_path: str, footer_text: str = "",
                              theme: dict = CTDI_THEME) -> str:
    blocks = _parse_blocks(markdown_text)
    slides = _blocks_to_slides(blocks)
    return render_pptx_branded(slides, output_path, footer_text=footer_text, theme=theme)


# ---------------------------------------------------------------------------
# PDF conversion -- LibreOffice headless. Works on the .docx/.pptx this
# module produces, or any other file format LibreOffice can open.
# ---------------------------------------------------------------------------

def to_pdf(path: str, timeout: int = 120) -> str:
    """Convert a docx/pptx/odt/etc. to PDF via `soffice --headless`.
    Returns the output PDF path (same directory, same basename, .pdf)."""
    src = Path(path)
    if not src.exists():
        raise FileNotFoundError(path)
    outdir = src.parent
    result = subprocess.run(
        ["soffice", "--headless", "--convert-to", "pdf", "--outdir", str(outdir), str(src)],
        capture_output=True, text=True, timeout=timeout,
    )
    pdf_path = outdir / (src.stem + ".pdf")
    if result.returncode != 0 or not pdf_path.exists():
        raise RuntimeError(
            f"soffice conversion failed (rc={result.returncode}): "
            f"{result.stdout}\n{result.stderr}"
        )
    return str(pdf_path)


def main() -> int:
    if len(sys.argv) < 3:
        print(__doc__)
        return 2
    mode = sys.argv[1]
    if mode == "pdf":
        print(to_pdf(sys.argv[2]))
        return 0
    if mode not in ("docx", "pptx"):
        print(f"unknown mode: {mode}", file=sys.stderr)
        return 2
    if len(sys.argv) < 4:
        print(f"usage: python3 -m second_brain.doc_generation {mode} input.md output.{mode}", file=sys.stderr)
        return 2
    input_path, output_path = sys.argv[2], sys.argv[3]
    text = Path(input_path).read_text()
    title = Path(input_path).stem
    if mode == "docx":
        print(markdown_to_docx(text, output_path, title=title))
    else:
        print(markdown_to_pptx(text, output_path, title=title))
    return 0


if __name__ == "__main__":
    sys.exit(main())
